#!/usr/bin/env python3
"""
Генератор презентаций DigitalAuditor Cleshnya
"""

import sys
from pathlib import Path
from datetime import datetime
import yaml
import subprocess
import shutil

class PresentationGenerator:
    def __init__(self, audience: str):
        self.audience = audience
        self.pres_dir = Path(__file__).parent
        self.slides_dir = self.pres_dir / "slides" / audience
        self.output_dir = self.pres_dir / "output"
        self.output_dir.mkdir(exist_ok=True)
        self.data = self._load_audit_data()
        
    def _load_audit_data(self) -> dict:
        data = {
            "date": datetime.now().strftime("%d.%m.%Y"),
            "version": "1.0",
            "metrics": {"total_time": "13 минут", "chunks": 1096, "findings": 5, "coverage": "100%"},
        }
        
        metrics_file = self.pres_dir / "data" / "metrics.yaml"
        if metrics_file.exists():
            with open(metrics_file, 'r', encoding='utf-8') as f:
                loaded = yaml.safe_load(f) or {}
                data["metrics"].update(loaded)
                
        return data
    
    def _replace_variables(self, content: str) -> str:
        metrics = self.data["metrics"]
        
        findings = metrics.get("findings", 0)
        if isinstance(findings, list):
            findings_count = len(findings)
        else:
            findings_count = findings
            
        content = content.replace("{{DATE}}", self.data["date"])
        content = content.replace("{{VERSION}}", self.data["version"])
        content = content.replace("{{AUDIT_TIME}}", str(metrics.get("total_time", "13 минут")))
        content = content.replace("{{CHUNKS_COUNT}}", str(metrics.get("chunks", 1096)))
        content = content.replace("{{FINDINGS_COUNT}}", str(findings_count))
        content = content.replace("{{COVERAGE}}", str(metrics.get("coverage", "100%")))
        
        return content
    
    def generate_markdown(self) -> Path:
        combined = []
        slide_files = sorted(self.slides_dir.glob("*.md"))
        
        for slide_file in slide_files:
            content = slide_file.read_text(encoding='utf-8')
            content = self._replace_variables(content)
            combined.append(content)
            combined.append("\n\n---\n\n")
        
        combined_md = self.output_dir / f"{self.audience}_combined.md"
        combined_md.write_text("".join(combined), encoding='utf-8')
        return combined_md
    
    def generate(self):
        print(f"\n[*] Генерация для: {self.audience}")
        md_file = self.generate_markdown()
        print(f"[+] Markdown: {md_file}")
        
        # Попытка создать PPTX через python-pptx
        try:
            from pptx import Presentation
            prs = Presentation()
            
            content = md_file.read_text(encoding='utf-8')
            slides = content.split("\n\n---\n\n")
            
            for slide_text in slides:
                if not slide_text.strip():
                    continue
                
                lines = slide_text.strip().split("\n")
                if lines[0].startswith("# "):
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    title = slide.shapes.title
                    title.text = lines[0].replace("# ", "")
                    if len(slide.placeholders) > 1:
                        slide.placeholders[1].text = "\n".join(lines[1:4])
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    if lines and lines[0].startswith("## "):
                        slide.shapes.title.text = lines[0].replace("## ", "")
                        if len(slide.placeholders) > 1:
                            slide.placeholders[1].text = "\n".join(lines[1:])
            
            pptx_file = self.output_dir / f"{self.audience}_{self.data['date'].replace('.', '-')}.pptx"
            prs.save(str(pptx_file))
            print(f"[+] PPTX: {pptx_file}")
            
        except ImportError:
            print("[-] python-pptx не установлен")
            print("[*] Доступен только Markdown")
        except Exception as e:
            print(f"[-] Ошибка PPTX: {e}")
        
        return md_file

def main():
    if len(sys.argv) < 2:
        print("Использование: python generate.py [executive|practitioner|technical|all]")
        sys.exit(1)
    
    audience = sys.argv[1]
    
    if audience == "all":
        for aud in ["executive", "practitioner", "technical"]:
            gen = PresentationGenerator(aud)
            gen.generate()
    else:
        gen = PresentationGenerator(audience)
        gen.generate()
    
    print("\n[+] Готово!")

if __name__ == "__main__":
    main()
